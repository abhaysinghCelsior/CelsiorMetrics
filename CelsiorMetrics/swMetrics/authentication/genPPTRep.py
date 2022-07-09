from plotly.graph_objs import Bar
import plotly.graph_objects as go
from plotly import offline
import plotly.express as px
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from docx.shared import Pt
import time
import array as arr
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
#from configMet import REP_PATH
from pptx.enum.text import PP_ALIGN
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

from .models import WeightedDefectDensity,OnTimeDelivery,QARemovalEfficiency,AutomatedUnitTestingEffectiveness,WeightedCodeReviewEffectiveness,EffortVariance
from .models import MetricsBenchMark,DefectSource,RC,DefectSummary,MetricsMonth,AppMetrics,ProjectOwner
from swMetrics.authentication.configMet import EXT,REP_NM,REP_PATH,SMTP,FROM,PORT,EMAIL_BODY,EMAIL_TAIL,REP_NM,MET1_TITLE,MET2_TITLE,MET3_TITLE,MET4_TITLE,MET5_TITLE,MET6_TITLE,SlideNoStartingMet,NO_of_MET
import swMetrics.authentication.Generic as G


def ProduceRep(MON,Flg,Opt1,Opt2,VALID_MONTH,PPT_Name,PPT_TITLE):
    ppt = Presentation(REP_PATH.strip()+"METRICS-Template.pptx")
    print(REP_PATH.strip()+"METRICS-Template.pptx")
    SlideNum=1
    MET=['WDD','OTD','QRE','AUTE','WCRE','EV']
    TITLE=['WEIGHTED DEFECT DENSITY','ON TIME DELIVERY','QA REMOVAL EFFICIENCY','AUTOMATED UNIT TESTING EFFECTIVENESS','WEIGHTED CODE REVIEW EFFECTIVENESS','EFFORT VARIENCE']
    OBJ=[WeightedDefectDensity,OnTimeDelivery,QARemovalEfficiency,AutomatedUnitTestingEffectiveness,WeightedCodeReviewEffectiveness,EffortVariance]
    REMARK=[MET1_TITLE,MET2_TITLE,MET3_TITLE,MET4_TITLE,MET5_TITLE,MET6_TITLE]
    for slide in ppt.slides:
        if(SlideNum==1):

                left =  Inches(2.75)
                height =  Inches(1)
                top = Inches(3)
                width=Inches(10)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                # adding Paragraphs
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                # adding text
                p.text = PPT_TITLE
                p.font.bold = True
                p.font.name= 'Montserrat SemiBold'
                p.font.size = Pt(40)
        if((SlideNum>=SlideNoStartingMet) and (SlideNum<=SlideNoStartingMet+NO_of_MET-1)):
                ret,Remarks=genMetrics(MON,MET[SlideNum-2],TITLE[SlideNum-2],OBJ[SlideNum-2],Flg,Opt1,Opt2,VALID_MONTH)
                if ret==0:
                        return -1 # In Valid

                img=REP_PATH.strip()+MET[SlideNum-2]+".png"
                pic = slide.shapes.add_picture(img, left= Inches(1),top = Inches(1),height = Inches(4),width = Inches(11))
                left = Inches(1)
                height = Inches(1.5)
                width=Inches(11)
                top=Inches(5)
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(220, 229, 204)
                top=Inches(4.6)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                # adding text
                p.text = REMARK[SlideNum-2]
                p.font.bold = True
                p.font.name= 'Montserrat SemiBold'
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(8, 28, 71)
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                # adding text
                if(len(Remarks.strip())>0):
                        p.text = "Outliers Reason:"+Remarks
                        p.font.color.rgb = RGBColor(245, 13, 6)
                else:
                        p.text = "All Data Points are within Benchmark"
                        p.font.color.rgb = RGBColor(13, 117, 27)
                p.font.bold = True
                p.font.name= 'Montserrat SemiBold'
                p.font.size = Pt(10)
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
        if(SlideNum==SlideNoStartingMet+NO_of_MET):
                genDefSummary(MON)
                img=REP_PATH.strip()+'DefectSummary.png'
                pic = slide.shapes.add_picture(img, left= Inches(0),top = Inches(1),height = Inches(6),width = Inches(13))
        SlideNum=SlideNum+1
    ppt.save(REP_PATH.strip()+PPT_Name)
    PRJ_DAFAULTER=G.genListProjDidntAddedMetricsData()
    sndMail("METRICS GENERATED",EMAIL_BODY+PRJ_DAFAULTER+EMAIL_TAIL,REP_PATH.strip()+PPT_Name)
    return 0

def getPPTInfo(Flg,Opt1,Opt2):
    VALID_MONTH=[]
    FIRST=""
    FrstFlg=False
    for key in MetricsMonth.objects.all().order_by('-Month'):
        if(FrstFlg is False):
            FIRST=  key.Month
            FrstFlg=True
        if G.ifMonthInOptedPriod(key.Month,Flg,Opt1,Opt2,FIRST):
            VALID_MONTH.append(key.Month)
                
    PPT_Name,PPT_TITLE=getFileName(Flg,Opt1,Opt2,VALID_MONTH,FIRST)
    VALID_MONTH.reverse()
    
    return PPT_Name,PPT_TITLE,VALID_MONTH
        
def genMetrics(MON,MET,TITLE,OBJ,Flg,Opt1,Opt2,VALID_MONTH):
        MET_I=[]
        BNCHMRK=[]
        KEY=[]
        BNCHMRK=[]
        MON=[]
        PRJ=[]

        #VALID_MONTH.reverse()
        DataFormat=".2%"
        YDataFormat="%{y}"
        Remarks=""
        print("-------------genMetrics:MON-->PRJ-->BENCHMARK--------"+MET)
        MET_DP=0.0
        for ValMon in VALID_MONTH:
                for key in OBJ.objects.filter(Month_of_Metrics=ValMon):
                    P=key.Project_Key.split("_")
                    KEY.append(key.Month_of_Metrics+":"+P[2])
                    OprToChk='G'
                    if(MET=="OTD"):
                            OprToChk='L'
                            MET_DP=key.OTD
                    elif (MET=="EV"):
                            MET_DP=key.EV
                    elif (MET=="AUTE"):
                            OprToChk='L'
                            DataFormat=".2"
                            MET_DP=key.AUTE
                    elif (MET=="QRE"):
                            OprToChk='L'
                            MET_DP=key.DRE
                    elif (MET=="WDD"):
                            DataFormat=".3"
                            MET_DP=key.WDD
                    elif (MET=="WCRE"):
                            OprToChk='L'
                            DataFormat=".2"
                            MET_DP=key.WCRE
                    else:
                            return 0,Remarks
                    MET_I.append(MET_DP)  
                    BnMrk_DP= MetricsBenchMark.objects.get(Metrics=MET,Month=ValMon)
                    MON.append('`'+ValMon)
                    PRJ.append(P[2])
                    BNCHMRK.append(BnMrk_DP.BenchMark)
                    if(OprToChk=='G'):
                            if(MET_DP>BnMrk_DP.BenchMark):
                                    if (key.Remarks is not None):
                                        Remarks=Remarks+ValMon+"-"+P[2]+":"+key.Remarks+"\n"
                    else:
                            if(MET_DP<BnMrk_DP.BenchMark):
                                    if (key.Remarks is not None):
                                            Remarks=Remarks+ValMon+"-"+P[2]+":"+key.Remarks+"\n"

                           
        data={'KEY':KEY,'Month':MON,"Project":PRJ,'BNCHMRK':BNCHMRK,MET+'_Achieved':MET_I}
        df = pd.DataFrame(data)
        print(df)
        fig = px.bar(df, x = 'Month', y = MET+'_Achieved', color = 'Project',barmode = 'group',labels={"Project": "Project"})
        fig.update_traces(texttemplate=YDataFormat)
        fig.update_layout(yaxis_tickformat=DataFormat)

        fig.add_traces(go.Scatter(x= df.Month, y=df.BNCHMRK, mode = "lines+text",text=df.BNCHMRK, name='Benchmark'))
        fig.update_traces(texttemplate=YDataFormat)
        fig.update_layout(yaxis_tickformat=DataFormat)
        fig.update_xaxes(showline=True, linewidth=1, linecolor='#040a60', mirror=True,ticks='outside')
        fig.update_yaxes(showline=True, linewidth=1, linecolor='#040a60', mirror=True,ticks='outside')
   
        # strip down the rest of the plot
        fig.update_layout(
            showlegend=True,
            paper_bgcolor="#cbb2fb",
            plot_bgcolor="#FBE0FB",
            margin=dict(t=10,l=10,b=10,r=10),
            xaxis=dict(
                title='Month',
                titlefont_size=24,
                tickfont_size=14,
                titlefont_color='#0a5c04',
                tickfont_color='#5c0504',
                       ),
            yaxis=dict(
                title=MET,
                titlefont_size=24,
                tickfont_size=14,
                titlefont_color='#0a5c04',
                tickfont_color= '#5c0504',
            ),

            legend=dict(
                x=1.01,
                y=1.01,
                bgcolor='rgba(255, 255, 255, 0)',
                bordercolor='#6FBD6A',
                ),
            bargap=0.1,
        )
     
        fname=REP_PATH.strip()+MET+'.png'
        fig.write_image(fname)
        return 1,Remarks

def genDefSummary(MON):
        Hdr=['Project', 'Critical#','High#','Medium#','Low#','Defect Cause','Source Phase of Defect','Remarks']
        try:
                DefSum_L= DefectSummary.objects.filter(Month_of_Metrics=MON).order_by('Month_of_Metrics','Project_Key')
        except DefectSummary.DoesNotExist:
                DefSum_L=None
        if DefSum_L is None:
                return 'E'

        C={}
        H={}
        M={}
        L={}
        RC=[]
        DS=[]
        RCC={}
        DSS={}
        Last_Prj=""
        first=1
        CC=[]
        HH=[]
        MM=[]
        LL=[]
        PRJ1=[]
        RC_CHOICES,DEFECTSOURCE_CHOICES=G.getRCandDS("N")
        Remark=[]
        print("-----------before for loop-----------")
        print(DefSum_L)
        for DefSum_Itr in DefSum_L:
                P=DefSum_Itr.Project_Key.split("-")
                PRJ=P[2] 
                if(DefSum_Itr.Severity=='Critical'):
                        C[PRJ]= 0
                if(DefSum_Itr.Severity=='High'):
                        H[PRJ]= 0
                if(DefSum_Itr.Severity=='Medium'):
                        M[PRJ]=0
                if(DefSum_Itr.Severity=='Low'):
                        L[PRJ]= 0

        
        for DefSum_Itr in DefSum_L:
                P=DefSum_Itr.Project_Key.split("-")
                PRJ=P[2]
                if(PRJ != Last_Prj):
                        if( first ==0 ):
                                str1=""
                                for idx in RCC.keys():
                                        if RCC[idx]>0:
                                                str1=str1+idx+":"+str(RCC[idx])+", "
                                RC.append(str1)
                                str1=""
                                for idx in DSS.keys():
                                        if DSS[idx]>0:
                                                str1=str1+idx+":"+str(DSS[idx])+", "
                                DS.append(str1)
                        
                        for idx in DEFECTSOURCE_CHOICES:
                                DSS[idx]=0
                        for idx in RC_CHOICES:
                                RCC[idx]=0
                        Remark=""
                
                RCC[DefSum_Itr.RC]=RCC[DefSum_Itr.RC]+1
                DSS[DefSum_Itr.DefectSource]=DSS[DefSum_Itr.DefectSource]+1
               
                if(DefSum_Itr.Severity=='Critical'):
                        C[PRJ]= C[PRJ]+1
                if(DefSum_Itr.Severity=='High'):
                        H[PRJ]= H[PRJ]+1
                if(DefSum_Itr.Severity=='Medium'):
                        M[PRJ]= M[PRJ]+1
                if(DefSum_Itr.Severity=='Low'):
                        L[PRJ]= L[PRJ]+1

                first=0
                Last_Prj=PRJ     


        str1=""
        
        for idx in RCC.keys():
                if RCC[idx]>0:
                        str1=str1+idx+":"+str(RCC[idx])+", "
        RC.append(str1)
        str1=""
        for idx in DSS.keys():
                if DSS[idx]>0:
                        str1=str1+idx+":"+str(DSS[idx])+","
        DS.append(str1) 
        
        for key in C.keys():
                PRJ1.append(key)
                CC.append(C[key])
                HH.append(H[key])
                MM.append(M[key])
                LL.append(L[key])
        
        fig = go.Figure(data=[go.Table(header=dict(values=Hdr),
                         cells=dict(values=[PRJ1,CC,HH,MM,LL,RC,DS, Remark]))
                             ])
        fig.data[0]['columnwidth'] = [8, 4,4,4,4,15,15,15]
        fig.update_layout(margin=dict(l=2, r=2, t=2, b=2),autosize=False,font=dict(family="Montserrat SemiBold",size=9))
        fname=REP_PATH.strip()+'DefectSummary.png'
        fig.write_image(fname)

def sndMail(SUBJECT,BODY,fileToAttach):
        msg = MIMEMultipart()
        msg['From'] = FROM
        msg['To'] = FROM
        msg['Subject'] = SUBJECT
        body = BODY
        msg.attach(MIMEText(body, 'plain'))
        attachment = open(fileToAttach, "rb")
        p = MIMEBase('application', 'octet-stream')
        p.set_payload((attachment).read())
        encoders.encode_base64(p)
        p.add_header('Content-Disposition', "attachment; filename= %s" % fileToAttach)
        msg.attach(p)
        s = smtplib.SMTP(SMTP, PORT)
        s.starttls()
        s.login(FROM, "NewWorld@2022")
        text = msg.as_string()
        s.sendmail(FROM, FROM, text)
        s.quit()

def valReport(request,YR,FR_L,TO_L):
    opt=request.POST.get('Period', "None")
    print(opt)
    if opt != "None":
        return "1",opt,None

    YRS=request.POST.get('YR', "None")
    print(YRS)
    if YRS != "None":
        for OP in YR:
            if(OP!=YRS):
                YR.remove(OP)
        return "2",YRS,None

    FR=request.POST.get('FR', "None")
    TO=request.POST.get('TO', "None")
    print(":"+FR+": -:"+TO+":")
    if ((FR == "None") and (TO == "None")):
        ret="E4"
    else:
        if((FR != "None") and (TO == "None")):
                ret="E2"
        elif((FR == "None") and (TO != "None")):
                ret="E3"
        else:
                if((FR != "None") and (TO != "None")):
                        if (TO <= FR):
                                ret= "E1"
                        else:
                                ret= "3",FR,TO
                else:
                        ret="E5"

    if(ret=='3'):
        if FR is not None:
          for OP in FR_L :
                if(OP!=FR):
                        FR_L.remove(OP)
        if TO is not None:
          for OP in TO_L :
                if(OP!=TO):
                        TO_L.remove(OP)
   
        return "3",FR,TO
    else:
        return ret,None,None

def getFileName(Flg,Opt1,Opt2,VALID_MONTH,FIRST):
        FN=REP_NM+"_"
        TITLE=""
        print(Opt1)
        
        if (Flg=="1"):    
                FR_MN=G.MonthNumToStr(G.MonthStrToNum(FIRST)-int(Opt1)+1)
                TITLE=FR_MN + " TO " + FIRST
                FN=FN+FR_MN+"_"+FIRST
        if(Flg=="2"):
                TITLE="Year:" + Opt1
                FN=FN+Opt1
        if(Flg=="3"):
                TITLE=Opt1 + " TO " + Opt2
                FN=FN+Opt1+"_"+Opt2
        FN=FN+EXT
        return FN, TITLE
